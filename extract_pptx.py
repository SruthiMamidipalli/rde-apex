from pptx import Presentation
import os

pptx_path = r'apex_loyalty_case_study.pptx'
prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides, 1):
    print("")
    print("=" * 80)
    print(f"SLIDE {i}")
    print("=" * 80)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(text)
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                print(row_text)
